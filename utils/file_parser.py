import PyPDF2
import docx
from pptx import Presentation

def extract_text_from_file(filepath):
    if filepath.endswith('.pdf'):
        return extract_pdf(filepath)
    elif filepath.endswith('.docx'):
        return extract_docx(filepath)
    elif filepath.endswith('.pptx'):
        return extract_pptx(filepath)
    return ""

def extract_pdf(filepath):
    text = ''
    with open(filepath, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            text += page.extract_text() or ''
    return text

def extract_docx(filepath):
    doc = docx.Document(filepath)
    return '\n'.join([para.text for para in doc.paragraphs])

def extract_pptx(filepath):
    prs = Presentation(filepath)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + '\n'
    return text
